__author__ = 'morrj140'

import logging
from nl_lib import Logger
from nl_lib.Concepts import Concepts
from nl_lib.ConceptGraph import PatternGraph, NetworkXGraph, Neo4JGraph, GraphVizGraph
from nl_lib.Constants import *

logger = Logger.setupLogging(__name__)

import math

from pptx import Presentation
from lxml import etree

from traceback import format_exc

import ImportArchi as ia

# Constants
EMU = 914400.0

gdb = "http://localhost:7474/db/data/"
#gdb = "http://10.92.82.60:7574/db/data/"

def addGraphNodes(graph, concepts, n=0):
    n += 1
    for c in concepts.getConcepts().values():

        logger.debug("%s[%d]" % (c.name, len(c.name)))

        if len(c.name.strip(" ")) == 0:
            return

        if not (c.typeName in ("Source", "Target")):
            return

        logger.debug("%d : %d Node c : %s:%s" % (n, len(c.getConcepts()), c.name, c.typeName))

        graph.addConcept(c)
        if len(c.getConcepts()) != 0:
            addGraphNodes(graph, c, n)

def addGraphEdges(graph, concepts, n=0):
    n += 1
    i = 1
    for c in concepts.getConcepts().values():

        if (c.name in ("l", "h", "t", "w")):
            return

        logger.debug("%d : %d Edge c : %s:%s" % (n, len(c.getConcepts()), c.name, c.typeName))
        if i == 1:
            p = c
            i += 1
        else:
            graph.addEdge(p, c)
        if len(c.getConcepts()) != 0:
            addGraphEdges(graph, c, n)

def graphConcepts(concepts, graph=None):

    if graph == None:
        #graph = Neo4JGraph(gdb)

        #logger.info("Clear the Graph @" + gdb)
        #graph.clearGraphDB()

        graph = GraphVizGraph()
        #graph = NetworkXGraph()
        #graph = PatternGraph()

    logger.info("Adding nodes the graph ...")
    addGraphNodes(graph, concepts)
    logger.info("Adding edges the graph ...")
    addGraphEdges(graph, concepts)

    if isinstance(graph, GraphVizGraph):
        filename="example.png"
        graph.exportGraph(filename=filename)
        logger.info("Saved Graph - %s" % filename)

    if isinstance(graph, Neo4JGraph):
        graph.setNodeLabels()

    if isinstance(graph, NetworkXGraph):
        graph.drawGraph("concepts.png")

        filename = "concepts.net"
        logger.info("Saving Pajek - %s" % filename)
        graph.saveGraphPajek(filename)

        graph.saveGraph("concepts.gml")
        logger.info("Saving Graph - %s" % "concepts.gml")

    if isinstance(graph, PatternGraph):
        logger.info("Exporting Graph")
        graph.exportGraph()

def findID(nid, dictNodes):
    try:
        for x in dictNodes.keys():
            logger.debug("    dictNodes[%s] : %s" % (dictNodes[x], x))

            if nid in dictNodes[x]:
                logger.debug("Found %s in %s" % (x, dictNodes[x]))
                return x
    except:
        em = format_exc().split('\n')[-2]
        logger.warn("findID : Warning: %s" % (em))

    return None

def findXY(nid, d):

    ld = list()

    try:
        ld = d[nid]

        logger.debug("ld : %s" % ld)

    except:
        pass

    return ld

def logList(l, n=0):

    n += 1
    s = " " * n

    logger.info("%sn=%d" % (s, n))

    for x in l:
        #logger.info("%sx=%s" % (s, x))
        if isinstance(x, list):
            logger.info("%slist: %s" % (s, x))
            logList(x, n)
        elif isinstance(x, dict):
            logger.info("%sdict: %s" % (s, x))
            logList(x, n)
        elif isinstance(x, tuple):
            logger.info("%stuple: %s" % (s, x))
            logList(x, n)
        else:
            if isinstance(x, str):
                logger.info("%sstr: %s" % (s, x))
            elif isinstance(x, float):
                logger.info("%sfloat: %3.2f" % (s, x))
            elif isinstance(x, int):
                logger.info("%sint: %d" % (s, x))


def shapeText(shape):
    name = ""
    if shape.has_text_frame:
        text_frame = shape.text_frame

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                logger.debug("%s" % run.text)
                name = name + run.text + " "

    return name

def shapeDim(shape, dictNodeXY):
    t = shape.top / EMU
    l = shape.left / EMU
    h = shape.height / EMU
    w = shape.width / EMU

    nid = shape.id

    dictDim = dict()
    dictDim["t"] = t
    dictDim["l"] = l
    dictDim["h"] = h
    dictDim["w"] = w
    dictNodeXY[nid] = dictDim

    logger.debug("shape.top     : %3.2f" % (t))
    logger.debug("shape.left    : %3.2f" % (l))
    logger.debug("shape.height  : %3.2f" % (h))
    logger.debug("shape.width   : %3.2f" % (w))
    logger.debug("shape.shape_type    : %s" % shape.shape_type)

    return nid, t, l, h , w

def addDictNodes(nid, name, dictNodes):

    name = ia.cleanString(name.rstrip(" ").lstrip(" "))

    if not (len(name) > 0):
        logger.warn("No Name!")
        return

    if dictNodes.has_key(name):
        nl = dictNodes[name]
        nl.append(nid)
        logger.debug("Duplicate Keys %s...%s" % (name, dictNodes[name]))
    else:
        nl = list()
        nl.append(nid)
        dictNodes[name] = nl


def addDictEdges(nid, xl, dictEdges):
    nxl = list()
    for x in xl:
        nxl.append(int(x))
        logger.debug ("%d:%s" % (nid,x))

    lenNXL = len(nxl)

    #
    # Only add connections between two nodes
    #
    if lenNXL == 3:
        if dictEdges.has_key(nid):
            nl = dictEdges[nid]
            nl.append(nxl)
            logger.debug("Duplicate Edges ...%s" % (dictEdges[nid]))
        else:
            el = list()
            el.append(nxl)
            dictEdges[nid] = el
    else:
        logger.debug("Only %d Connectors!" % (len(nxl)))

    return lenNXL

def showConcepts(concepts):
    n = 0
    for x in concepts.getConcepts().values():
        n += 1
        logger.info("x %s[%s]" % (x.name, x.typeName))
        for y in x.getConcepts().values():
            logger.info("  y %s[%s]" % (y.name, y.typeName))
            for z in y.getConcepts().values():
                if not (z.name in ("h", "l", "t", "w")):
                    logger.info("    z  %s[%s]" % (z.name, z.typeName))

def getPoint(d):

    t = d["t"]
    l = d["l"]
    h = d["h"]
    w = d["w"]

    py = t + ( h / 2.0)
    px = l + ( h / 2.0)

    return px, py

def lineMagnitude (x1, y1, x2, y2):
    lineMagnitude = math.sqrt(math.pow((x2 - x1), 2)+ math.pow((y2 - y1), 2))
    return lineMagnitude

#Calc minimum distance from a point and a line segment (i.e. consecutive vertices in a polyline).
def DistancePointLine (px, py, x1, y1, x2, y2):
    try:
        #http://local.wasp.uwa.edu.au/~pbourke/geometry/pointline/source.vba
        LineMag = lineMagnitude(x1, y1, x2, y2)

        u1 = (((px - x1) * (x2 - x1)) + ((py - y1) * (y2 - y1)))
        u = u1 / (LineMag * LineMag)

        if (u < 0.00001) or (u > 1):
            #// closest point does not fall within the line segment, take the shorter distance
            #// to an endpoint
            ix = lineMagnitude(px, py, x1, y1)
            iy = lineMagnitude(px, py, x2, y2)
            if ix > iy:
                DistancePointLine = iy
            else:
                DistancePointLine = ix
        else:
            # Intersecting point is on the line, use the formula
            ix = x1 + u * (x2 - x1)
            iy = y1 + u * (y2 - y1)
            DistancePointLine = lineMagnitude(px, py, ix, iy)

        return DistancePointLine
    except:
        return 0

def crawlPPTX(concepts, path_to_presentation):

    prs = Presentation(path_to_presentation)

    sNum = 0

    for slide in prs.slides:
        dictNodes = dict()
        dictEdges = dict()
        dictText  = dict()

        dictNodeXY = dict()
        dictTextXY = dict()

        logger.debug ("--new slide--")
        logger.debug("%s" % slide.partname)
        logger.debug("slideName : %s" % slide.name)

        sNum += 1

        #
        # Get Title of Slide
        #
        titleSlide = ""

        for idx, ph in enumerate(slide.shapes.placeholders):
            logger.debug ("**    %s:%s    **" % (idx, ph.text))
            if idx == 0:
                titleSlide = ph.text

        u = ia.cleanString(titleSlide)

        logger.info("%d.%s" % (sNum, u))
        tss = "%d.%s" % (sNum, u)
        q = concepts.addConceptKeyType(tss, "Slide")

        #showConcepts(concepts)

        #
        # Iterate ihrough slides
        #
        n = 0
        nc = 0
        for shape in slide.shapes:
            logger.debug ("...%s..." % type(shape))
            logger.debug("shape.element.xml : %s" % shape.element.xml)
            logger.debug("shape.name : %s[%d]" % (shape.name,  shape.id - 1))

            n += 1

            sn = shape.name

            nid = shape.id

            # Get Shape Info
            if shape.name[:5] in ("Recta", "Round", "Strai"):
                nid, t, l, h, w = shapeDim(shape, dictNodeXY)

                tl = (l, t)
                tr = (l + w, t)
                bl = (l, t + h)
                br = (l + w, t + h)

                name = shapeText(shape)

                if len(name) > 1:
                    logger.info("  node : %s[%d] - %s" % (name, nid, shape.name))

                    addDictNodes(nid, name, dictNodes)

                    b = q.addConceptKeyType(ia.cleanString(name), "Node")
                    b.addConceptKeyType("t", str(t))
                    b.addConceptKeyType("l", str(l))
                    b.addConceptKeyType("h", str(h))
                    b.addConceptKeyType("w", str(w))

            #
            # Add in Connections
            #
            elif sn.find("Connector") != -1:

                xmlShape = shape.element.xml

                logger.debug("xmlShape : %s" % xmlShape)

                tree = etree.fromstring(xmlShape)

                xl = tree.xpath("//@id")

                logger.debug("xl : %s" % xl)

                numEdges = addDictEdges(nid, xl, dictEdges)

                if numEdges == 3:
                    nc += 1
                    logger.info("  %d Found Edge %d" % (nc, shape.id))


            #
            # Get Text boxes and associate with Connector
            #
            elif shape.name[:8] in ("Text Box", "TextBox "):

                nid, t, l, h, w = shapeDim(shape, dictTextXY)

                name = shapeText(shape)

                if name != None:
                    nxl = list()
                    nxl.append(nid)
                    dictText[name] = nxl

                    logger.info("  TextBox : %s[%d]" % (name, shape.id))

            else:
                logger.debug("Skipped : %s" % shape.name)

        #
        # Now match the Connector with text
        #

        listEdges = dictEdges.values()

        logger.info("listEdges : %d" % len(listEdges))

        tbFound = 0
        tbTotal = len(dictTextXY)
        logger.info("Search for %s Text Box Connector's" % len(dictTextXY))

        for txt in dictTextXY.keys():
            logger.debug("txt : %s[%s]" % (txt, dictTextXY[txt] ))
            searchText = findID(txt, dictText)

            logger.info("  Search Text : %s" % (searchText))

            # get text point - middle of node
            px, py = getPoint(dictTextXY[txt])

            cDist = 1000.0
            cNode = None
            csn = None
            ctn = None

            # for each node in dictEdges
            ni = 0
            for edge in listEdges:

                logger.debug("  edge: %s" % edge)

                try:
                    # get source
                    source = edge[0][2]
                    sName = findID(source, dictNodes)
                    sl = dictNodeXY[source]
                    spx, spy = getPoint(sl)

                    # get target
                    target = edge[0][1]
                    tName = findID(target, dictNodes)
                    tl = dictNodeXY[target]

                    tpx, tpy = getPoint(tl)

                    # determine distance between points
                    d = DistancePointLine (px, py, spx, spy, tpx, tpy)

                    if d < cDist:
                        cDist = d
                        cNode = edge[0][0]
                        csn = sName
                        tsn = tName

                except:
                    pass

            if cNode != None:
                tbFound += 1
                logger.debug("    Closest Connector : %s" % cNode)
                logger.info("    found(%d:%d] - %s->%s->%s [%2.3f]" % (tbFound, tbTotal, csn, searchText, tsn, cDist))

                edge   = searchText
                source = sName
                target = tName

                dimSource = sl
                dimTarget = tl

                if edge == None:
                    edge = "TBD"

                d = q.getConcepts()[csn]

                for ld in dimSource.keys():
                    logger.debug("%s %s:%2.3f" % (source, ld, dimSource[ld]))
                    d.addConceptKeyType(ld, str(dimSource[ld]))

                    f = d.addConceptKeyType(target, "Target")
                    for ld in dimTarget.keys():
                        logger.debug("%s %s:%2.3f" % (target, ld, dimSource[ld]))
                        f.addConceptKeyType(ld, str(dimTarget[ld]))

                    f.addConceptKeyType(ia.cleanString(edge), "Edge")


        if tbTotal != 0:
                logger.info("Found [%3.1f] Text Box Connectors" % ((tbFound / float(tbTotal)) * 100.0))

        dictTextXY = dict()

if __name__ == "__main__":
    #graph = NetworkXGraph()

    #path_to_presentation = "/Users/morrj140/PycharmProjects/ArchiConcepts/example2.pptx"
    #path_to_presentation = "/Users/morrj140/PycharmProjects/ArchiConcepts/ARP-TBX - High Level Solution_Draft_v9.pptx"
    #path_to_presentation = "/Users/morrj140/Development/GitRepository/ArchiConcepts/ARP-TBX - High Level Solution_Draft_v10a.pptx"
    #path_to_presentation = "/Users/morrj140/Development/GitRepository/ArchiConcepts/Accovia_Replacement_Messages.pptx"
    #path_to_presentation = "/Users/morrj140/Development/GitRepository/ArchiConcepts/ARP-TBX - Next Level Solution -v1.pptx"

    path_to_presentation = "/Users/morrj140/Development/GitRepository/ArchiConcepts/simple2.pptx"

    c = Concepts("Application", "Relations")

    crawlPPTX(c, path_to_presentation)

    #graphConcepts(c)

    c.logConcepts()

    Concepts.saveConcepts(c, "pptx.p")


