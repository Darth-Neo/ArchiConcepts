#!/usr/bin/python
#
# Create PPTX from Archimate XML
#
__author__ = u'morrj140'
__VERSION__ = u'0.1'

import sys
import os
import StringIO
import glob, time, math, zipfile

from Logger import *
logger = setupLogging(__name__)
logger.setLevel(INFO)

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
#from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
#from pptx.oxml.shapes import connector
#from pptx.parts.slide import _SlideShapeTree

from ArchiLib import ArchiLib
from Constants import *

import pytest

class ArchiCreatePPTX(object):
    filePPTXIn  = None
    filePPTXOut = None

    def __init__(self, fileArchimate, filePPTXIn, filePPTXOut):
        self.A_NS           =  u"http://schemas.openxmlformats.org/drawingml/2006/main"
        self.P_NS           =  u"http://schemas.openxmlformats.org/presentationml/2006/main"
        self.R_NS           =  u"http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        self.namespacesPPTX = {u"p" : self.P_NS, u"a" : self.A_NS, u"r" : self.R_NS}

        self.SLD_LAYOUT_TITLE_AND_CONTENT = 1
        self.TITLE_ONLY_SLIDE_LAYOUT = 5

        self.SCALE = 0.90
        self.EMU = 914400.0

        self.filePPTXIn     = filePPTXIn
        self.filePPTXOut    = filePPTXOut
        self.fileArchimate  = fileArchimate

        if os.path.isfile(self.fileArchimate) <> True:
            logger.error(u"File does not exist : %s" % self.fileArchimate)

        etree.QName(ARCHIMATE_NS, u'model')
        self.tree = etree.parse(self.fileArchimate)

        self.prs = Presentation()

    # Example of what the xml for a connector looks like
    def addXMLConnector(self, shape):
        name = u"Straight Arrow Connector 43"
        id = 34
        sourceID = 21
        targetID = 9

        t = shape.top / self.EMU
        l = shape.left / self.EMU
        h = shape.height / self.EMU
        w = shape.width / self.EMU

        nid = shape.id
        shape.name = u"Straight Arrow Connector 43"

        logger.debug(u"shape.top     : %3.2f" % (t))
        logger.debug(u"shape.left    : %3.2f" % (l))
        logger.debug(u"shape.height  : %3.2f" % (h))
        logger.debug(u"shape.width   : %3.2f" % (w))
        logger.debug(u"shape.shape_type    : %s" % shape.shape_type)
        xmlConnector = u" \
             <p:cxnSp xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\" " \
                u"xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\" " \
                u"xmlns:r=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships\"> \
              <p:nvCxnSpPr> \
                <p:cNvPr id=\"34\" name=\"Straight Arrow Connector 43\"/> \
                <p:cNvCxnSpPr> \
                  <a:stCxn id=\"21\" idx=\"3\"/> \
                  <a:endCxn id=\"9\" idx=\"1\"/> \
                </p:cNvCxnSpPr> \
                <p:nvPr/> \
              </p:nvCxnSpPr> \
              <p:spPr> \
                <a:xfrm flipV=\"1\">  \
                  <a:off x=\"5056479\" y=\"3507277\"/>  \
                  <a:ext cx=\"1642504\" cy=\"12357\"/> \
                </a:xfrm> \
                <a:prstGeom prst=\"straightConnector1\"> <a:avLst/></a:prstGeom> \
                <a:ln> <a:headEnd type=\"none\"/> <a:tailEnd type=\"arrow\"/> </a:ln> \
              </p:spPr> \
              <p:style> \
                <a:lnRef idx=\"1\"> <a:schemeClr val=\"dk1\"/> </a:lnRef>  \
                <a:fillRef idx=\"0\"> <a:schemeClr val=\"dk1\"/> </a:fillRef> \
                <a:effectRef idx=\"0\"> <a:schemeClr val=\"dk1\"/> </a:effectRef>  \
                <a:fontRef idx=\"minor\"> \
                    <a:schemeClr val=\"tx1\"/>  \
                </a:fontRef> \
              </p:style>  \
            </p:cxnSp> "

        shape.element.xml = xmlConnector

    def getNode(self, n, listModels, type):

        attributes = n.attrib

        if attributes.get(ARCHI_TYPE) == type:
            if attributes.get(u"id") is not None:
                listModels.append((n, attributes))

                logger.debug(u"%s : %s:%s:%s:%s" % (DIAGRAM_MODEL, n.tag, n.get(u"name"), n.get(u"id"), attributes.get(ARCHI_TYPE)))

        for y in n:
            self.getNode(y, listModels, type)

    def getAll(self, listModels, type=DIAGRAM_MODEL):
        for x in self.tree.getroot():
            self.getNode(x, listModels, type)

    def findNode(self, id, tag=u"element"):
        logger.debug(u"id = %s" % id)
        xp = u"//%s[@id='%s']" % (tag, id)
        stp = self.tree.xpath(xp)

        if len(stp) > 0:
            return stp[0]

        return stp

    def findDiagramObject(self, listDO, value):
        for x in listDO:
            if x[0] == u"archimate:DiagramObject":
                logger.debug(u"value : %s[%s] x : %s[%s]" % (value, type(value), x[3], type(x[3])))
                if x[3] == value:
                    logger.debug(u"Found!")
                    return x

        return None

    def project(self, x, y, scale=None):

        if scale is None:
            scale = self.SCALE
        try:
            x = (((float(x) / 100.0)) * scale)
            y = (((float(y) / 100.0)) * scale)
            return x, y
        except:
            return None, None

    def findShape(self, sm, listDO):

        for model in listDO:
            if model == sm:
                return model[6]
        return None

    def fixSlides(self, listSlides):

        # unzip .pptx to temporary space and remove file
        timestamp = unicode(time.time()).replace(u".", u"")
        zip_file = zipfile.ZipFile(self.PPTXFilename, u"r")
        zip_file.extractall(os.path.join(u"./tmp", timestamp))
        zip_file.close

        os.remove(self.PPTXFilename)

        # register necessary namespaces
        etree.register_namespace(u"a", u"http://schemas.openxmlformats.org/drawingml/2006/main")
        etree.register_namespace(u"p", u"http://schemas.openxmlformats.org/presentationml/2006/main")
        etree.register_namespace(u"r", u"http://schemas.openxmlformats.org/officeDocument/2006/relationships")

        # parse xml document and find shape tree
        slideNum = 1
        slideFilename = u"slide%d.xml" % slideNum
        tree = etree.parse(os.path.join(u"./tmp", timestamp, u"ppt", u"slides", slideFilename))

        root = tree.getroot()

        max_sp = 0
        xp = u"//@id"
        for sp in tree.xpath(xp):
            if int(sp) > max_sp:
                max_sp = int(sp)
            logger.debug(u"sp : %s" % (sp))

        n = 1
        for connector in self.listConnectors:

            connectorID = int(max_sp) + n
            n += 1

            start = connector[0][6]
            start_left   = connector[0][7]
            start_top    = connector[0][8]
            start_width  = connector[0][9]
            start_height = connector[0][10]

            logger.debug(u"StartID : %s[%s]" % (start, connector[0][2]))
            logger.debug(u"    l:%d,t:%d,w:%d,h:%d)" %(start_left, start_top, start_width, start_height))

            end = connector[1][6]
            end_left   = connector[1][7]
            end_top    = connector[1][8]
            end_width  = connector[1][9]
            end_height = connector[1][10]

            logger.debug(u"EndID : %s[%s]" % (end, connector[1][2]))
            logger.debug(u"    l:%d,t:%d,w:%d,h:%d)" %(end_left, end_top, end_width, end_height))

            sxml_id = int(start)
            start_idx = int(start) + 1
            exml_id = int(end)
            end_idx  = int(end + 1)

            cxn_x = start_left
            cxn_y = start_top

            cxn_cx = end_left
            cxn_cy = end_top

            connectionShape = etree.SubElement(root, u"{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp", nsmap=self.namespacesPPTX)

            nonVisualConnectorShapeDrawingProperties = etree.Element(u"{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr")
            cNonVisualProperties = etree.SubElement(nonVisualConnectorShapeDrawingProperties, u"{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
            cNonVisualProperties.set(u"id", unicode(connectorID + 1))
            cNonVisualProperties.set(u"name", u"Straight Arrow Connector " + unicode(connectorID))
            cNonVisualConnectorShapeDrawingProperties = etree.SubElement(nonVisualConnectorShapeDrawingProperties, u"{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr")
            etree.SubElement(nonVisualConnectorShapeDrawingProperties, u"{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr")

            connectionStart = etree.SubElement(cNonVisualConnectorShapeDrawingProperties, u"{http://schemas.openxmlformats.org/drawingml/2006/main}stCxn")
            connectionStart.set(u"id", unicode(sxml_id))     # shape index from which connector starts (param)
            connectionStart.set(u"idx", unicode(start_idx))  # connector spawn point index
            connectionEnd = etree.SubElement(cNonVisualConnectorShapeDrawingProperties, u"{http://schemas.openxmlformats.org/drawingml/2006/main}endCxn")
            connectionEnd.set(u"id", unicode(exml_id))       # shape index at which connector ends (param)
            connectionEnd.set(u"idx", unicode(end_idx))      # connector termination point index

            shapeProperties = etree.SubElement(connectionShape, u"{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
            twodTransform = etree.SubElement(shapeProperties, u"{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")

            twodTransform.set(u"flipH", u"1")

            # location of bounding box (params)
            offset = etree.SubElement(twodTransform, u"{http://schemas.openxmlformats.org/drawingml/2006/main}off")
            offset.set(u"x", unicode(cxn_x).replace(u".0", u""))
            offset.set(u"y", unicode(cxn_y).replace(u".0", u""))

            # height and width of bounding box (params)
            extents = etree.SubElement(twodTransform, u"{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
            extents.set(u"cx", unicode(cxn_cx).replace(u".0", u""))
            extents.set(u"cy", unicode(cxn_cy).replace(u".0", u""))

            presetGeometry = etree.SubElement(shapeProperties, u"{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
            presetGeometry.set(u"prst", u"line")
            etree.SubElement(presetGeometry, u"{http://schemas.openxmlformats.org/drawingml/2006/main}avLst")

            style = etree.SubElement(connectionShape, u"{http://schemas.openxmlformats.org/presentationml/2006/main}style")

            lineReference = etree.SubElement(style, u"{http://schemas.openxmlformats.org/drawingml/2006/main}lnRef")
            lineReference.set(u"idx", u"1")
            schemeColor = etree.SubElement(lineReference, u"{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr")
            schemeColor.set(u"val", u"dk1")

            fillReference = etree.SubElement(style, u"{http://schemas.openxmlformats.org/drawingml/2006/main}fillRef")
            fillReference.set(u"idx", u"0")
            schemeColor = etree.SubElement(fillReference, u"{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr")
            schemeColor.set(u"val", u"dk1")

            effectReference = etree.SubElement(style, u"{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef")
            effectReference.set(u"idx", u"0")
            schemeColor = etree.SubElement(effectReference, u"{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr")
            schemeColor.set(u"val", u"dk1")

            fontReference = etree.SubElement(style, u"{http://schemas.openxmlformats.org/drawingml/2006/main}fontRef")
            fontReference.set(u"idx", u"minor")
            schemeColor = etree.SubElement(fontReference, u"{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr")
            schemeColor.set(u"val", u"tx1")


        # write out the final tree
        f = open(os.path.join(u"./tmp", timestamp, u"ppt", u"slides", u"slide1.xml"), u"w")
        tree.write(f, encoding=u"UTF-8", xml_declaration=True)
        f.close()

        # zip the file
        zip_file = zipfile.ZipFile(self.PPTXFilename, u"w", zipfile.ZIP_DEFLATED)
        os.chdir(os.path.join(u"./tmp", timestamp))
        for root, dirs, files in os.walk(u"."):
            for file in files:
                zip_file.write(os.path.join(root, file))
        zip_file.close

    def buildPPTX(self):

        title_only_slide_layout = self.prs.slide_layouts[self.TITLE_ONLY_SLIDE_LAYOUT]

        slide = self.prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        timeTxt = time.strftime(u"%Y%d%m_%H%M%S")
        shapes.title.text = u"Built on %s" % (timeTxt)

        listModels = list()
        listConnectors = list()

        # From Archimate File
        self.getAll(listModels)

        #
        # Iterate through all Archimate Diagrams
        #
        logger.info(u"====Add Archimate Diagrams====")
        max_sp = 0

        for x in listModels:
            #
            # Since you must know the shape_id for each corresponding Diagram Object, make them first
            #
            listDO = list()
            logger.info(u"%s[%s]" % (x[0].get(u"name"), x[0].get(u"id")))

            slideName = unicode(x[0].get(u"name"))
            ls = list()
            ls.append(u"Slide Title")
            ls.append(slideName)
            listDO.append(ls)

            p = u"//element[@id=\"%s\"]" % (x[0].get(u"id"))
            r = self.tree.xpath(p, namespaces=NS_MAP)
            xc = r[0].getchildren()

            for y in xc:
                child = unicode(y.get(u"archimateElement"))
                logger.debug(u"  %s[%s]: entity:%s" % (y.get(ARCHI_TYPE), y.get(u"id"), child))

                n = self.findNode(child)

                if n == None or isinstance(n, list):
                    continue

                shapeName = unicode(n.get(u"name"))
                logger.debug(u"  DO = %s" % (shapeName))

                z = y.getchildren()
                for w in z:
                    logger.debug(u"    %s[%s]" % (w.get(ARCHI_TYPE), w.get(u"id")))

                    if w.get(ARCHI_TYPE) == u"archimate:Connection":
                        logger.debug(u"      source=%s, target=%s, relationship=%s" % (w.get(u"source"), w.get(u"target"), w.get(u"relationship")))

                        ls = list()
                        ls.append(u"Connector")
                        ls.append(unicode(w.get(u"source")))
                        ls.append(unicode(w.get(u"target")))

                        relation = w.get(u"relationship")
                        rn = self.findNode(relation)
                        logger.debug(u"      relation : %s[%s]" % (rn.get(u"name"), rn.get(ARCHI_TYPE)))
                        if rn.get(u"name") is not None:
                            ls.append(unicode(rn.get(u"name")))

                        listDO.append(ls)

                        sn = self.findNode(w.get(u"source"), tag=u"child")
                        logger.debug(u"  SO = %s" % sn.get(u"name"))

                        st = self.findNode(w.get(u"target"), tag=u"child")
                        logger.debug(u"  TO = %s" % st.get(u"name"))

                    elif w.get(u"x") is not None:
                        logger.debug(u"    x=%s, y=%s" % (w.get(u"x"), w.get(u"y")))
                        ls = list()
                        ls.append(unicode(y.get(ARCHI_TYPE)))
                        ls.append(unicode(n.get(ARCHI_TYPE)))
                        ls.append(shapeName)
                        ls.append(unicode(y.get(u"id")))
                        ls.append(unicode(w.get(u"x")))
                        ls.append(unicode(w.get(u"y")))
                        listDO.append(ls)

            #
            # Now iterate the the diagram objects and add shapes to Slide
            # plus grab the shape_id!
            #
            logger.debug(u"====Add Slide====")
            for model in listDO:
                if model[0] == u"Slide Title":
                    #
                    # New Slide
                    #
                    slide_layout = self.prs.slide_layouts[self.TITLE_ONLY_SLIDE_LAYOUT]
                    slide = self.prs.slides.add_slide(slide_layout)
                    shapes = slide.shapes

                    shapes.title.text = u"%s" % (model[1])

                    logger.debug(u"model : %s" % model[1])

                elif model[0] == u"archimate:DiagramObject":
                    logger.debug(u"%s" % model)
                    name = model[2]
                    x, y = self.project(model[4], model[5])

                    if x is None or y is None:
                        continue

                    left   = Inches(x)
                    top    = Inches(y)

                    width  = Inches(1.0 * self.SCALE)
                    height = Inches(0.75 * self.SCALE)

                    logger.debug(u"DiagramObject : %s(%s,%s)" % (model[2], model[4], model[5]))
                    logger.debug(u"    l:%d,t:%d,w:%d,h:%d)" %(left, top, width, height))
                    logger.debug(u"    Point (%d:%d)" % (left, top))

                    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

                    # capture shape id!
                    model.append(shape.id)
                    model.append(left)
                    model.append(top)
                    model.append(width)
                    model.append(height)

                    logger.debug(u"shape.id : %s" % shape.id)

                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = name

                    font = run.font
                    font.name = u"Calibri"
                    font.size = Pt(10 * self.SCALE)
                    font.color.rgb = RGBColor(50, 50, 50) # grey

                    # set shape fill
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(204, 224, 255)

                elif model[0] == u"Connector":
                    logger.debug(u"  model[0] %s, model[1] %s" % (model[1], model[2]))

                    source = self.findDiagramObject(listDO, model[1])
                    target = self.findDiagramObject(listDO, model[2])

                    logger.debug(u"  source %s" % (source))
                    logger.debug(u"  target %s" % (target))

                    ll = list()
                    ll.append(source)
                    ll.append(target)
                    listConnectors.append(ll)

        # save file
        logger.info(u"\n Saved %s" % self.filePPTXOut)
        self.prs.save(self.filePPTXOut)

        return self.prs

        #
        # Add Connectors
        #
        # logger.debug("====Add Connectors====")
        # fixSlides(PPTXFilename, listConnectors)

def test_ArchiCreatePPTX():

    start_time = ArchiLib.startTimer()

    createPPTX = ArchiCreatePPTX(fileArchimateTest, filePPTXIn, filePPTXOut)

    createPPTX.buildPPTX()

    ArchiLib.stopTimer(start_time)

if __name__ == u"__main__":
    test_ArchiCreatePPTX()