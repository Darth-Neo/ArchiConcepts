#!/usr/bin/python
#
# PPTX Crawl
#
__author__ = u'morrj140'
__VERSION__ = u'0.1'

import math

from Logger import *
logger = setupLogging(__name__)
logger.setLevel(INFO)

from nl_lib.Concepts import Concepts
from nl_lib.ConceptGraph import PatternGraph, NetworkXGraph, Neo4JGraph, GraphVizGraph
from nl_lib.Constants import *

from pptx import Presentation
from lxml import etree

from traceback import format_exc

from ArchiLib import ArchiLib
from Constants import *


class PPTXCreateArchil(object):

    graph      = None

    dictNodes  = None
    dictEdges  = None

    dictText   = None
    dictNodeXY = None
    dictTextXY = None

    def __init__(self, fileCrawl, fileArchimate):
        self.EMU = 914400.0
        self.fileArchimate = fileArchimate

        self.path_to_presentation = fileCrawl

        self.dictNodes = dict()
        self.dictEdges = dict()
        self.dictText  = dict()

        self.dictNodeXY = dict()
        self.dictTextXY = dict()

        self.al = ArchiLib(fileArchimate)

        self.graph = GraphVizGraph()
        # self.graph = NetworkXGraph()
        # self.graph = PatternGraph()

        self.prs = Presentation(self.path_to_presentation)

        self.concepts = Concepts(u"Application", u"Relations")


    def addGraphNodes(self, concepts, n=0):
        n += 1
        for c in concepts.getConcepts().values():

            logger.debug(u"%s[%d]" % (c.name, len(c.name)))

            if len(c.name.strip(u" ")) == 0:
                return

            if not (c.typeName in (u"Source", u"Target")):
                return

            logger.debug(u"%d : %d Node c : %s:%s" % (n, len(c.getConcepts()), c.name, c.typeName))

            self.graph.addConcept(c)
            if len(c.getConcepts()) != 0:
                self.addGraphNodes(c, n)

    def addGraphEdges(self, concepts, n=0):
        n += 1
        i = 1
        for c in concepts.getConcepts().values():

            if (c.name in (u"l", u"h", u"t", u"w")):
                return

            logger.debug(u"%d : %d Edge c : %s:%s" % (n, len(c.getConcepts()), c.name, c.typeName))
            if i == 1:
                p = c
                i += 1
            else:
                self.graph.addEdge(p, c)
            if len(c.getConcepts()) != 0:
                self.addGraphEdges(c, n)

    def graphConcepts(self, concepts, graph=None):

        logger.info(u"Adding nodes the graph ...")
        self.addGraphNodes(concepts)

        logger.info(u"Adding edges the graph ...")
        self.addGraphEdges(concepts)

        if isinstance(graph, GraphVizGraph):
            filename = u"example.png"
            graph.exportGraph(filename=filename)
            logger.info(u"Saved Graph - %s" % filename)

        if isinstance(graph, Neo4JGraph):
            graph.setNodeLabels()

        if isinstance(graph, NetworkXGraph):
            graph.drawGraph(u"concepts.png")

            filename = u"concepts.net"
            logger.info(u"Saving Pajek - %s" % filename)
            graph.saveGraphPajek(filename)

            graph.saveGraph(u"concepts.gml")
            logger.info(u"Saving Graph - %s" % u"concepts.gml")

        if isinstance(graph, PatternGraph):
            logger.info(u"Exporting Graph")
            graph.exportGraph()

    def findID(self, nid):
        try:
            for x in self.dictNodes.keys():
                logger.debug(u"    dictNodes[%s] : %s" % (self.dictNodes[x], x))

                if nid in self.dictNodes[x]:
                    logger.debug(u"Found %s in %s" % (x, self.dictNodes[x]))
                    return x
        except:
            em = format_exc().split('\n')[-2]
            logger.warn(u"findID : Warning: %s" % (em))

        return None

    def findXY(self, nid, d):

        ld = list()

        try:
            ld = d[nid]

            logger.debug(u"ld : %s" % ld)

        except:
            pass

        return ld

    def logList(self, l, n=0):

        n += 1
        s = " " * n

        logger.info(u"%sn=%d" % (s, n))

        for x in l:
            # logger.info("%sx=%s" % (s, x))
            if isinstance(x, list):
                logger.info(u"%slist: %s" % (s, x))
                self.logList(x, n)
            elif isinstance(x, dict):
                logger.info(u"%sdict: %s" % (s, x))
                self.logList(x, n)
            elif isinstance(x, tuple):
                logger.info(u"%stuple: %s" % (s, x))
                self.logList(x, n)
            else:
                if isinstance(x, str):
                    logger.info(u"%sstr: %s" % (s, x))
                elif isinstance(x, float):
                    logger.info(u"%sfloat: %3.2f" % (s, x))
                elif isinstance(x, int):
                    logger.info(u"%sint: %d" % (s, x))

    def shapeText(self, shape):
        name = u""
        if shape.has_text_frame:
            text_frame = shape.text_frame

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    logger.debug(u"%s" % run.text)
                    name = name + run.text + u" "

        return name

    def shapeDim(self, shape):
        t = shape.top / self.EMU
        l = shape.left / self.EMU
        h = shape.height / self.EMU
        w = shape.width / self.EMU

        nid = shape.id

        dictDim = dict()
        dictDim[u"t"] = t
        dictDim[u"l"] = l
        dictDim[u"h"] = h
        dictDim[u"w"] = w

        self.dictNodeXY[nid] = dictDim

        logger.debug(u"shape.top     : %3.2f" % (t))
        logger.debug(u"shape.left    : %3.2f" % (l))
        logger.debug(u"shape.height  : %3.2f" % (h))
        logger.debug(u"shape.width   : %3.2f" % (w))
        logger.debug(u"shape.shape_type    : %s" % shape.shape_type)

        return nid, t, l, h , w

    def addDictNodes(self, nid, name):

        name = unicode(name).rstrip(u" ").lstrip(u" ")

        if not (len(name) > 0):
            logger.warn(u"No Name!")
            return

        if name in self.dictNodes:
            nl = self.dictNodes[name]
            nl.append(nid)
            logger.debug(u"Duplicate Keys %s...%s" % (name, self.dictNodes[name]))
        else:
            nl = list()
            nl.append(nid)
            self.dictNodes[name] = nl

    def addDictEdges(self, nid, xl):
        nxl = list()
        for x in xl:
            nxl.append(int(x))
            logger.debug(u"%d:%s" % (nid, x))

        lenNXL = len(nxl)

        #
        # Only add connections between two nodes
        #
        if lenNXL == 3:
            if self.dictEdges.has_key(nid):
                nl = self.dictEdges[nid]
                nl.append(nxl)
                logger.debug(u"Duplicate Edges ...%s" % (self.dictEdges[nid]))
            else:
                el = list()
                el.append(nxl)
                self.dictEdges[nid] = el
        else:
            logger.debug(u"Only %d Connectors!" % (len(nxl)))

        return lenNXL

    def showConcepts(self, concepts):
        n = 0
        for x in concepts.getConcepts().values():
            n += 1
            logger.info(u"x %s[%s]" % (x.name, x.typeName))
            for y in x.getConcepts().values():
                logger.info(u"  y %s[%s]" % (y.name, y.typeName))
                for z in y.getConcepts().values():
                    if not (z.name in (u"h", u"l", u"t", u"w")):
                        logger.info(u"    z  %s[%s]" % (z.name, z.typeName))

    def getPoint(self, d):

        t = d[u"t"]
        l = d[u"l"]
        h = d[u"h"]
        w = d[u"w"]

        py = t + (h / 2.0)
        px = l + (h / 2.0)

        return px, py

    def lineMagnitude (self, x1, y1, x2, y2):
        lineMagnitude = math.sqrt(math.pow((x2 - x1), 2)+ math.pow((y2 - y1), 2))
        return lineMagnitude

    # Calc minimum distance from a point and a line segment (i.e. consecutive vertices in a polyline).
    def DistancePointLine (self, px, py, x1, y1, x2, y2):
        try:
            # http://local.wasp.uwa.edu.au/~pbourke/geometry/pointline/source.vba
            LineMag = self.lineMagnitude(x1, y1, x2, y2)

            u1 = (((px - x1) * (x2 - x1)) + ((py - y1) * (y2 - y1)))
            u = u1 / (LineMag * LineMag)

            if (u < 0.00001) or (u > 1):
                # closest point does not fall within the line segment, take the shorter distance
                # to an endpoint
                ix = self.lineMagnitude(px, py, x1, y1)
                iy = self.lineMagnitude(px, py, x2, y2)
                if ix > iy:
                    DistancePointLine = iy
                else:
                    DistancePointLine = ix
            else:
                # Intersecting point is on the line, use the formula
                ix = x1 + u * (x2 - x1)
                iy = y1 + u * (y2 - y1)
                DistancePointLine = self.lineMagnitude(px, py, ix, iy)

            return DistancePointLine
        except:
            return 0

    def crawlPPTX(self):

        sNum = 0

        for slide in self.prs.slides:

            logger.debug(u"--new slide--")
            logger.debug(u"%s" % slide.partname)
            logger.debug(u"slideName : %s" % slide.name)

            sNum += 1

            #
            # Get Title of Slide
            #
            titleSlide = u""

            for idx, ph in enumerate(slide.shapes.placeholders):
                # logger.debug ("**    %s:%s    **" % (idx, ph.text))
                if idx == 0:
                    titleSlide = ph.text

            u = self.al.cleanString(titleSlide)

            logger.info(u"%d.%s" % (sNum, u))
            tss = u"%d.%s" % (sNum, u)
            q = self.concepts.addConceptKeyType(tss, u"Slide")

            # showConcepts(concepts)

            #
            # Iterate ihrough slides
            #
            n = 0
            nc = 0
            for shape in slide.shapes:
                logger.debug(u"...%s..." % type(shape))
                logger.debug(u"shape.element.xml : %s" % shape.element.xml)
                logger.debug(u"shape.name : %s[%d]" % (shape.name,  shape.id - 1))

                n += 1

                sn = shape.name

                nid = shape.id

                # Get Shape Info
                if shape.name[:5] in (u"Recta", u"Round", u"Strai"):
                    nid, t, l, h, w = self.shapeDim(shape)

                    tl = (l, t)
                    tr = (l + w, t)
                    bl = (l, t + h)
                    br = (l + w, t + h)

                    name = self.shapeText(shape)

                    if len(name) > 1:
                        logger.info(u"  node : %s[%d] - %s" % (name, nid, shape.name))

                        self.addDictNodes(nid, name)

                        b = q.addConceptKeyType(self.al.cleanString(name), u"Node")
                        b.addConceptKeyType(u"t", str(t))
                        b.addConceptKeyType(u"l", str(l))
                        b.addConceptKeyType(u"h", str(h))
                        b.addConceptKeyType(u"w", str(w))

                #
                # Add in Connections
                #
                elif sn.find(u"Connector") != -1:

                    xmlShape = shape.element.xml

                    logger.debug(u"xmlShape : %s" % xmlShape)

                    tree = etree.fromstring(xmlShape)

                    xl = tree.xpath(u"//@id")

                    logger.debug(u"xl : %s" % xl)

                    numEdges = self.addDictEdges(nid, xl)

                    if numEdges == 3:
                        nc += 1
                        logger.info(u"  %d Found Edge %d" % (nc, shape.id))

                #
                # Get Text boxes and associate with Connector
                #
                elif shape.name[:8] in (u"Text Box", u"TextBox "):

                    nid, t, l, h, w = self.shapeDim(shape)

                    name = self.shapeText(shape)

                    if name is not None:
                        nxl = list()
                        nxl.append(nid)
                        self.dictText[name] = nxl

                        logger.info(u"  TextBox : %s[%d]" % (name, shape.id))

                else:
                    logger.debug(u"Skipped : %s" % shape.name)

            #
            # Now match the Connector with text
            #

            listEdges = self.dictEdges.values()

            logger.info(u"listEdges : %d" % len(listEdges))

            tbFound = 0
            tbTotal = len(self.dictTextXY)
            logger.info(u"Search for %s Text Box Connector's" % len(self.dictTextXY))

            for txt in self.dictTextXY.keys():
                logger.debug(u"txt : %s[%s]" % (txt, dictTextXY[txt] ))
                searchText = self.findID(txt, self.dictText)

                logger.info(u"  Search Text : %s" % (searchText))

                # get text point - middle of node
                px, py = self.getPoint(dictTextXY[txt])

                cDist = 1000.0
                cNode = None
                csn = None
                ctn = None

                # for each node in dictEdges
                ni = 0
                for edge in listEdges:

                    logger.debug(u"  edge: %s" % edge)

                    try:
                        # get source
                        source = edge[0][2]
                        sName = self.findID(source)
                        sl = self.dictNodeXY[source]
                        spx, spy = self.getPoint(sl)

                        # get target
                        target = edge[0][1]
                        tName = self.findID(target)
                        tl = self.dictNodeXY[target]

                        tpx, tpy = self.getPoint(tl)

                        # determine distance between points
                        d = self.DistancePointLine (px, py, spx, spy, tpx, tpy)

                        if d < cDist:
                            cDist = d
                            cNode = edge[0][0]
                            csn = sName
                            tsn = tName

                    except:
                        pass

                if cNode != None:
                    tbFound += 1
                    logger.debug(u"    Closest Connector : %s" % cNode)
                    logger.info(u"    found(%d:%d] - %s->%s->%s [%2.3f]" % (tbFound, tbTotal, csn, searchText, tsn, cDist))

                    edge   = searchText
                    source = sName
                    target = tName

                    dimSource = sl
                    dimTarget = tl

                    if edge is None:
                        edge = u"TBD"

                    d = q.getConcepts()[csn]

                    for ld in dimSource.keys():
                        logger.debug(u"%s %s:%2.3f" % (source, ld, dimSource[ld]))
                        d.addConceptKeyType(ld, str(dimSource[ld]))

                        f = d.addConceptKeyType(target, u"Target")
                        for ld in dimTarget.keys():
                            logger.debug(u"%s %s:%2.3f" % (target, ld, dimSource[ld]))
                            f.addConceptKeyType(ld, str(dimTarget[ld]))

                        f.addConceptKeyType(self.al.cleanString(edge), u"Edge")


            if tbTotal != 0:
                    logger.info(u"Found [%3.1f] Text Box Connectors" % ((tbFound / float(tbTotal)) * 100.0))

            dictTextXY = dict()

        return self.concepts

def PPTXCreateArchi():

    start_time = ArchiLib.startTimer()

    logger.info(u"Using : %s" % filePPTXIn)

    cpptx = PPTXCreateArchil(filePPTXIn, fileArchimateTest)

    c = cpptx.crawlPPTX()

    c.logConcepts()

    Concepts.saveConcepts(c, fileConceptsPPTX)

    ArchiLib.stopTimer(start_time)

if __name__ == u"__main__":
    PPTXCreateArchi()