__author__ = 'morrj140'

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO
from pptx.util import Inches
SLD_LAYOUT_TITLE_AND_CONTENT = 1
TITLE_ONLY_SLIDE_LAYOUT = 5

#
# Title Slide
#
prs = Presentation()
title_only_slide_layout = prs.slide_layouts[TITLE_ONLY_SLIDE_LAYOUT]

slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding an AutoShape'

left = Inches(0.93)  # 0.93" centers this overall set of shapes
top = Inches(3.0)
width = Inches(1.75)
height = Inches(1.0)

#
# Adding Shapes
#

shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = 'Step 1'

left = left + width - Inches(0.4)
width = Inches(2.0)  # chevrons need more width for visual balance

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = 'Step %d' % n
    left = left + width - Inches(0.4)

#
# New Slide
#
slide_layout = prs.slide_layouts[TITLE_ONLY_SLIDE_LAYOUT]
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

left   = Inches(6.0)
top    = Inches(4.0)
width  = Inches(1.0)
height = Inches(1.0)

shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)

left   = Inches(2.0)
top    = Inches(2.0)
width  = Inches(1.0)
height = Inches(1.0)

shape = shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)

left = Inches(2.0)
top = Inches(3.0)
width = Inches(4)
height = Inches(2)

line = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left, top, width, height)
prs.save('test3.pptx')