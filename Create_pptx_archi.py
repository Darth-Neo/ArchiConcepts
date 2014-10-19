__author__ = 'morrj140'

import sys
import os
import StringIO
import glob, time, math, zipfile
from nl_lib import Logger
logger = Logger.setupLogging(__name__)
from nl_lib.Constants import *
from nl_lib.Concepts import Concepts
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
import import_artifacts as ia

SLD_LAYOUT_TITLE_AND_CONTENT = 1
TITLE_ONLY_SLIDE_LAYOUT = 5

namespaces={'xsi': 'http://www.w3.org/2001/XMLSchema-instance', 'archimate': 'http://www.archimatetool.com/archimate'}
XML_NS         =  "http://www.w3.org/2001/XMLSchema-instance"
ARCHIMATE_NS   =  "http://www.archimatetool.com/archimate"
NS_MAP = {"xsi": XML_NS, "archimate" : ARCHIMATE_NS}
ARCHI_TYPE = "{http://www.w3.org/2001/XMLSchema-instance}type"

PPTXFilename = 'test3.pptx'
DIAGRAM_MODEL = "archimate:ArchimateDiagramModel"
SCALE = 0.90
EMU = 914400.0

# Example of what the xml for a connector looks like
def addXMLConnector(tree, shape):
    name = "Straight Arrow Connector 43"
    id = 34
    sourceID = 21
    targetID = 9

    t = shape.top / EMU
    l = shape.left / EMU
    h = shape.height / EMU
    w = shape.width / EMU

    nid = shape.id
    shape.name = "Straight Arrow Connector 43"

    logger.info("shape.top     : %3.2f" % (t))
    logger.info("shape.left    : %3.2f" % (l))
    logger.info("shape.height  : %3.2f" % (h))
    logger.info("shape.width   : %3.2f" % (w))
    logger.info("shape.shape_type    : %s" % shape.shape_type)
    xmlConnector = " \
         <p:cxnSp xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\" " \
            "xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\" " \
            "xmlns:r=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships\"> \
          <p:nvCxnSpPr> \
            <p:cNvPr id=\"34\" name=\"Straight Arrow Connector 43\"/> \
            <p:cNvCxnSpPr> \
              <a:stCxn id=\"21\" idx=\"3\"/> \
              <a:endCxn id=\"9\" idx=\"1\"/> \
            </p:cNvCxnSpPr> \
            <p:nvPr/> \
          </p:nvCxnSpPr> \
          <p:spPr> \
            <a:xfrm flipV=\"1\">  \
              <a:off x=\"5056479\" y=\"3507277\"/>  \
              <a:ext cx=\"1642504\" cy=\"12357\"/> \
            </a:xfrm> \
            <a:prstGeom prst=\"straightConnector1\"> <a:avLst/></a:prstGeom> \
            <a:ln> <a:headEnd type=\"none\"/> <a:tailEnd type=\"arrow\"/> </a:ln> \
          </p:spPr> \
          <p:style> \
            <a:lnRef idx=\"1\"> <a:schemeClr val=\"dk1\"/> </a:lnRef>  \
            <a:fillRef idx=\"0\"> <a:schemeClr val=\"dk1\"/> </a:fillRef> \
            <a:effectRef idx=\"0\"> <a:schemeClr val=\"dk1\"/> </a:effectRef>  \
            <a:fontRef idx=\"minor\"> \
                <a:schemeClr val=\"tx1\"/>  \
            </a:fontRef> \
          </p:style>  \
        </p:cxnSp> "

    shape.element.xml = xmlConnector

def getNode(n, listModels, type):

    attributes = n.attrib

    if attributes.get(ARCHI_TYPE) == type:
        if attributes.get("id") != None:
            listModels.append((n, attributes))

            logger.info("%s : %s:%s:%s:%s" % (DIAGRAM_MODEL, n.tag, n.get("name"), n.get("id"), attributes.get(ARCHI_TYPE)))

    for y in n:
        getNode(y, listModels, type)

def getAll(tree, listModels, type=DIAGRAM_MODEL):
    for x in tree.getroot():
        getNode(x, listModels, type)

def findNode(tree, id, tag="element"):
    logger.debug("id = %s" % id)
    xp = "//%s[@id='%s']" % (tag, id)
    stp = tree.xpath(xp)

    if len(stp) > 0:
        return stp[0]

    return stp

def findDiagramObject(listDO, value):
    for x in listDO:
        if x[0] == "archimate:DiagramObject":
            logger.debug("value : %s[%s] x : %s[%s]" % (value, type(value), x[3], type(x[3])))
            if x[3] == value:
                logger.debug("Found!")
                return x

    return None

def project(x, y, scale=SCALE):
    try:
        x = (((float(x) / 100.0)) * scale)
        y = (((float(y) / 100.0)) * scale)
        return x, y
    except:
        return None, None

def findShape(sm, listDO):

    for model in listDO:
        if model == sm:
            return model[6]

    return None

if __name__ == "__main__":
    filePPTX = "Archimate.pptx"
    fileArchimateIn = "/Users/morrj140/Documents/SolutionEngineering/Archimate Models/CodeGen_v23.archimate"
    etree.QName(ARCHIMATE_NS, 'model')
    tree = etree.parse(fileArchimateIn)

    ia.logAll(tree, "archimate:ApplicationComponent")

    #
    # Title Slide
    #

    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[TITLE_ONLY_SLIDE_LAYOUT]

    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    timeTxt = time.strftime("%Y%d%m_%H%M%S")
    shapes.title.text = "Built from %s on %s" % (fileArchimateIn, timeTxt)

    listModels = list()
    listSlides = list()

    getAll(tree, listModels)

    for x in listModels[:1]:
        listDO = list()
        logger.info("%s[%s]" % (x[0].get("name"), x[0].get("id")))

        slideName = str(x[0].get("name"))
        ls = list()
        ls.append("Slide Title")
        ls.append(slideName)
        listDO.append(ls)

        p = "//element[@id=\"%s\"]" % (x[0].get("id"))
        r = tree.xpath(p, namespaces=namespaces)
        xc = r[0].getchildren()

        for y in xc:
            child = str(y.get("archimateElement"))
            logger.info("  %s[%s]: entity:%s" % (y.get(ARCHI_TYPE), y.get("id"), child))

            n = findNode(tree, child)

            if n == None or isinstance(n, list):
                continue

            shapeName = str(n.get("name"))
            logger.info("  DO = %s" % (shapeName))

            z = y.getchildren()
            for w in z:
                logger.debug("    %s[%s]" % (w.get(ARCHI_TYPE), w.get("id")))

                if w.get(ARCHI_TYPE) == "archimate:Connection":
                    logger.info("      source=%s, target=%s" % (w.get("source"), w.get("target")))

                    ls = list()
                    ls.append("Connector")
                    ls.append(str(w.get("source")))
                    ls.append(str(w.get("target")))
                    listDO.append(ls)

                    sn = findNode(tree, w.get("source"), tag="child")
                    logger.debug("  SO = %s" % sn.get("name"))

                    st = findNode(tree, w.get("target"), tag="child")
                    logger.debug("  TO = %s" % st.get("name"))

                elif w.get("x") != None:
                    logger.info("    x=%s, y=%s" % (w.get("x"), w.get("y")))
                    ls = list()
                    ls.append(str(y.get(ARCHI_TYPE)))
                    ls.append(str(n.get(ARCHI_TYPE)))
                    ls.append(shapeName)
                    ls.append(str(y.get("id")))
                    ls.append(str(w.get("x")))
                    ls.append(str(w.get("y")))
                    listDO.append(ls)

        #
        # Add Presentation Slide
        #
        logger.info("====Add Slides====")
        for model in listDO:
            if model[0] == "Slide Title":
                #
                # New Slide
                #
                slide_layout = prs.slide_layouts[TITLE_ONLY_SLIDE_LAYOUT]
                slide = prs.slides.add_slide(slide_layout)
                shapes = slide.shapes

                shapes.title.text = "%s" % (model[1])

                logger.info("model : %s" % model[1])

            elif model[0] == "archimate:DiagramObject":
                logger.info("%s" % model)
                name = model[2]
                x, y = project(model[4], model[5])

                if x == None or y == None:
                    continue

                left   = Inches(x)
                top    = Inches(y)

                width  = Inches(1.0)
                height = Inches(0.75)

                logger.info("DiagramObject : %s(%s,%s)" % (model[2], model[4], model[5]))
                logger.info("    l:%d,t:%d,w:%d,h:%d)" %(left, top, width, height))
                logger.info("    Point (%d:%d)" % (left, top))

                shape = shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
                )

                # capture shape id
                model.append(shape.id)
                logger.info("shape.id : %s" % shape.id)

                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = name

                font = run.font
                font.name = "Calibri"
                font.size = Pt(10)
                font.color.rgb = RGBColor(50, 50, 50) # grey

                # set shape fill
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(204, 224, 255)

            elif model[0] == "Connector":
                logger.info("model[0] %s, model[1] %s" % (model[1], model[2]))

                source = findDiagramObject(listDO, model[1])
                target = findDiagramObject(listDO, model[2])

                logger.info("  source %s" % (source))
                logger.info("  target %s" % (target))

                ll = list()
                ll.append(source)
                ll.append(target)
                listSlides.append(ll)

    # save file
    prs.save(PPTXFilename)










